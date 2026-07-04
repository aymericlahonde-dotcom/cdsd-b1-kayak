"""Génère Presentation_Kayak.pptx — Bloc 1 Jedha (5 min jury).

Slide 1  : Titre + attentes Kayak
Slide 2  : Architecture pipeline (texte + flux)
Slide 3  : Collecte des données (chiffres clés)
Slide 4  : Score météo + Top 5 destinations
Slide 5  : Stockage AWS S3 + RDS PostgreSQL
Slide 6  : Conclusion + perspectives prod

Style étudiant : sobre, lisible, un peu technique, pas IA-perfect.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Palette Jedha-ish
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ORANGE = RGBColor(0xE8, 0x7A, 0x1E)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_title_bar(slide, title, subtitle=None):
    """Bandeau titre en haut, fond navy."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.0))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = ORANGE


def add_text_box(slide, left, top, width, height, lines, size=16, color=NAVY,
                 bold_first=False, bullet=True):
    """Boîte de texte simple avec une liste de lignes."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "• " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(6)
    return box


def add_notes(slide, notes_text):
    """Speaker notes (mode présentateur)."""
    slide.notes_slide.notes_text_frame.text = notes_text


def add_footer(slide, idx, total=6):
    """Petit footer avec numéro de slide + auteur."""
    f = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    tf = f.text_frame
    p = tf.paragraphs[0]
    p.text = f"Aymeric Lahonde — CDSD Jedha 2026 — Bloc 1 Kayak"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p2 = tf.add_paragraph()
    p2.text = ""
    # numéro à droite
    n = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35))
    pn = n.text_frame.paragraphs[0]
    pn.text = f"{idx} / {total}"
    pn.font.size = Pt(10)
    pn.font.color.rgb = GRAY
    pn.alignment = PP_ALIGN.RIGHT


# ============================================================
# SLIDE 1 — Titre + attentes
# ============================================================
s = prs.slides.add_slide(BLANK)
# Fond clair
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
# Bloc navy à gauche pour le titre
left_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), Inches(0.4), Inches(4.0))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = ORANGE; left_bar.line.fill.background()

# Titre principal
t = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(1.4))
p = t.text_frame.paragraphs[0]
p.text = "Plan your trip with Kayak"
p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = NAVY

t2 = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(12.0), Inches(0.6))
p = t2.text_frame.paragraphs[0]
p.text = "Bloc 1 — Construction et alimentation d'une base de données (Data Engineering / ETL)"
p.font.size = Pt(20); p.font.color.rgb = GRAY; p.font.italic = True

# Bloc "Ce qu'attend Kayak"
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.7),
                         Inches(11.7), Inches(2.7))
box.fill.solid(); box.fill.fore_color.rgb = WHITE
box.line.color.rgb = NAVY; box.line.width = Pt(1.5)
tf = box.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "Ce que Kayak attend"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ORANGE
for line in [
    "70% des utilisateurs veulent plus d'infos sur leur destination",
    "Recommander les meilleures villes françaises selon météo + hôtels",
    "Construire le pipeline de A à Z : pas de données existantes",
    "Restitution sous forme de cartes interactives (top 5 villes + top 20 hôtels)",
]:
    p = tf.add_paragraph()
    p.text = "•  " + line
    p.font.size = Pt(15); p.font.color.rgb = NAVY
    p.space_before = Pt(4)

# Auteur
auth = s.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(12.0), Inches(0.5))
p = auth.text_frame.paragraphs[0]
p.text = "Aymeric Lahonde · Promotion CDSD Jedha 2026 · RNCP35288"
p.font.size = Pt(13); p.font.color.rgb = GRAY

add_notes(s, """[~30 s — INTRO]

Bonjour, je suis Aymeric Lahonde, et je vais vous présenter le projet Bloc 1 de la certification Jedha CDSD : « Plan your trip with Kayak ».

Le contexte : Kayak, c'est un moteur de recherche de voyages qui appartient au groupe Booking Holdings. Une étude utilisateur leur a montré que 70% des gens qui planifient un voyage veulent plus d'infos sur leur destination. Du coup, ils veulent recommander les meilleures villes selon deux critères : la météo et les hôtels disponibles.

Le périmètre qu'on m'a confié : les 35 plus belles villes françaises. Et le défi, c'est qu'il n'y a PAS de données — il faut tout construire : la collecte, le stockage cloud, et la restitution.

Mon livrable : un pipeline complet, une base SQL en RDS, et 2 cartes interactives. Je vais vous montrer comment je m'y suis pris.""")

add_footer(s, 1)


# ============================================================
# SLIDE 2 — Architecture détaillée
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
add_title_bar(s, "Architecture du pipeline ETL",
              "Collecte → Data Lake → Data Warehouse → Visualisation")

# Schéma horizontal des 6 étapes
boxes_data = [
    ("Nominatim\n(geocoding)",  "Lat / Lon\n35 villes",      ORANGE),
    ("OpenWeather\nOne Call API", "Météo 5j\n(JSON)",         ORANGE),
    ("Booking.com\nSelenium + BS4", "Hôtels\n(scraping)",     ORANGE),
    ("AWS S3\nData Lake",         "CSV bruts\nprefix raw/",   NAVY),
    ("AWS RDS\nPostgreSQL",       "Tables\nnettoyées",        NAVY),
    ("Plotly\nMap viz",           "2 cartes\nHTML",           GREEN),
]
n = len(boxes_data)
total_w = Inches(12.6)
gap = Inches(0.18)
box_w = (total_w - gap * (n - 1)) / n
top = Inches(1.5)
left0 = Inches(0.35)
for i, (label, sub, col) in enumerate(boxes_data):
    left = left0 + (box_w + gap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(10); p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    # Flèche entre les boîtes
    if i < n - 1:
        ax = left + box_w + Emu(20000)
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax,
                                top + Inches(0.55), gap - Emu(40000), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = GRAY
        ar.line.fill.background()

# Stack technique en dessous
add_text_box(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(3.5),
             ["Stack technique",
              "Python 3.11 · pandas · requests · beautifulsoup4",
              "Selenium (Chrome headless) — Booking dynamique en JS",
              "boto3 → S3 · psycopg2 + SQLAlchemy → RDS",
              "Plotly Express → cartes interactives HTML"],
             size=15, bold_first=True, bullet=False)

add_text_box(s, Inches(7.0), Inches(3.4), Inches(6.0), Inches(3.5),
             ["Pourquoi cette architecture",
              "S3 = stockage brut, schéma libre, archivage pas cher",
              "RDS = données nettoyées, requêtable en SQL",
              "Découplage collecte / analyse → on peut rejouer",
              "Free tier AWS → 0 € pendant la formation"],
             size=15, bold_first=True, bullet=False)

add_notes(s, """[~1 min — ARCHITECTURE]

Voici l'architecture complète. On lit de gauche à droite, c'est un pipeline ETL classique en 6 étapes.

D'abord la COLLECTE :
- Nominatim me donne les coordonnées GPS des 35 villes — c'est gratuit, sans clé API, parfait pour un projet bootcamp.
- OpenWeather One Call API : pour chaque ville, je récupère la météo sur 5 jours en JSON.
- Booking.com : je scrape les hôtels avec Selenium parce que le site est dynamique — un simple requests.get ne récupère pas le contenu, le HTML est généré côté client par JavaScript.

Ensuite le STOCKAGE en deux niveaux — c'est important :
- AWS S3 = mon Data Lake. J'y dépose les CSV bruts. C'est ma source de vérité, peu chère, schéma libre.
- AWS RDS PostgreSQL = mon Data Warehouse. Je relis les CSV depuis S3, je nettoie, je charge dans des tables SQL.

Pourquoi les deux ? Parce que si je n'avais que RDS, je perds la donnée brute si l'ETL casse. Et si je n'avais que S3, je n'ai pas de SQL rapide pour les analystes. C'est complémentaire.

Enfin la RESTITUTION : Plotly Express me génère 2 cartes HTML interactives.

Côté stack : Python, Selenium pour le dynamique, boto3 pour S3, psycopg2 pour Postgres. Tout est en free tier AWS — donc 0 € pendant la formation.""")

add_footer(s, 2)


# ============================================================
# SLIDE 3 — Collecte des données
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
add_title_bar(s, "Collecte des données",
              "Géocodage · Météo OpenWeather · Scraping Booking.com")

# 3 grosses tuiles chiffres
tiles = [
    ("35 / 35", "villes géocodées",      "Nominatim — 1 req/sec\npolitique d'usage"),
    ("175",     "mesures météo",         "5 jours × 35 villes\nOne Call API"),
    ("679",     "hôtels scrapés",        "≈ 20 / ville\nSelenium + BeautifulSoup"),
]
for i, (big, mid, sub) in enumerate(tiles):
    left = Inches(0.5 + i * 4.25)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5),
                             Inches(4.0), Inches(2.4))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = NAVY; box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = big
    p.font.size = Pt(46); p.font.bold = True; p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = mid
    p.font.size = Pt(16); p.font.color.rgb = NAVY; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = sub
    p.font.size = Pt(12); p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)

# Bloc défis techniques
add_text_box(s, Inches(0.5), Inches(4.3), Inches(12.5), Inches(2.5),
             ["Défis techniques rencontrés",
              "Booking en JS dynamique → Selenium headless obligatoire (~10 s / ville)",
              "Nominatim limite 1 req/sec → time.sleep(1) entre chaque ville",
              "Sélecteurs CSS Booking changent souvent → code fragile, à monitorer",
              "Encoding accents (Saint-Malo → St Malo) → normalisation côté géocodage",
              "Try / except autour des appels API → on ne perd pas la collecte sur 1 ville"],
             size=15, bold_first=True, bullet=False)

add_notes(s, """[~1 min — COLLECTE]

Voici les chiffres bruts de la collecte.

35 villes sur 35, j'ai tout géocodé avec Nominatim. À 1 requête/seconde, ça prend 35 secondes. Pas de rate limit dépassé.

175 mesures météo : 5 jours × 35 villes. Une mesure par jour avec température max, probabilité de pluie, volume de pluie, humidité, vent. Tout ça vient de l'API One Call d'OpenWeather.

679 hôtels scrapés sur Booking, environ 20 par ville. C'est la partie la plus longue du pipeline — environ 10 secondes par ville en Selenium.

Côté défis techniques : Booking est un site dynamique, donc requests + BeautifulSoup ne suffit pas. Il faut Selenium qui pilote un vrai Chrome headless pour exécuter le JavaScript.

Le code n'est pas robuste à 100 % — si Booking change ses sélecteurs CSS, il faut mettre à jour le scraper. C'est inhérent au scraping, on l'accepte. J'ai mis des try/except partout pour qu'une ville qui plante ne perde pas la collecte des 34 autres.""")

add_footer(s, 3)


# ============================================================
# SLIDE 4 — Score météo + Top 5
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
add_title_bar(s, "Score météo & Top 5 destinations",
              "Métrique composite normalisée 0–100")

# Formule à gauche
add_text_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.5),
             ["Formule de scoring"], size=18, bold_first=True, bullet=False, color=ORANGE)

formula = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0),
                             Inches(6.0), Inches(2.4))
formula.fill.solid(); formula.fill.fore_color.rgb = WHITE
formula.line.color.rgb = NAVY; formula.line.width = Pt(1.0)
tf = formula.text_frame
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "score = 0.4 × T_max"
p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAVY
for line, color in [
    ("        + 0.3 × (1 − P_pluie)", NAVY),
    ("        + 0.2 × (1 − V_pluie)", NAVY),
    ("        + 0.1 × (1 − Humidité)", NAVY),
    ("", WHITE),
    ("Toutes les composantes normalisées 0–100", GRAY),
    ("Pondération arbitraire mais documentée", GRAY),
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14) if "Pondération" not in line and "Toutes" not in line else Pt(12)
    p.font.color.rgb = color
    p.font.italic = "Pondération" in line or "Toutes" in line

# Top 5 à droite — tableau visuel
add_text_box(s, Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5),
             ["Top 5 destinations (score / 100)"], size=18, bold_first=True,
             bullet=False, color=ORANGE)

top5 = [
    ("1", "Aix-en-Provence", "88,9"),
    ("2", "Avignon",         "87,8"),
    ("3", "Marseille",       "86,7"),
    ("4", "Paris",           "85,5"),
    ("5", "Grenoble",        "82,3"),
]
for i, (rank, city, score) in enumerate(top5):
    top = Inches(2.0 + i * 0.55)
    row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), top,
                             Inches(6.0), Inches(0.5))
    row.fill.solid()
    row.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG
    row.line.color.rgb = GRAY; row.line.width = Pt(0.5)
    tf = row.text_frame
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = f"  {rank}.   {city:<25}            {score}"
    p.font.size = Pt(15); p.font.color.rgb = NAVY
    if i == 0:
        p.font.bold = True
        p.font.color.rgb = ORANGE

# Note explicative en bas
add_text_box(s, Inches(0.5), Inches(5.5), Inches(12.5), Inches(1.5),
             ["Lecture : le score privilégie les villes chaudes et sèches sur la fenêtre 5 jours.",
              "Dominante Sud / Sud-Est cohérente avec un printemps méditerranéen.",
              "Limite assumée : pas de saisonnalité, pas de prix hôtels — voir « pour aller plus loin »."],
             size=13, bullet=True, color=GRAY)

add_notes(s, """[~1 min — SCORE & TOP 5]

Maintenant le cœur métier : comment je classe les villes.

J'ai construit un score composite sur 100, avec 4 composantes pondérées :
- 40 % pour la température max — c'est ce qui compte le plus en vacances
- 30 % pour la probabilité de pluie inversée — moins de pluie = meilleur score
- 20 % pour le volume de pluie inversé
- 10 % pour l'humidité inversée

Toutes les composantes sont normalisées entre 0 et 100. La pondération est arbitraire — je l'assume — mais elle est documentée et explicable. En prod chez Kayak, on apprendrait ces poids à partir des comportements réels : clics, réservations, temps passé sur la page.

Le résultat : le Top 5 c'est Aix-en-Provence à 88,9, Avignon, Marseille, Paris, Grenoble. C'est très Sud-Est, ce qui est cohérent avec une fenêtre de 5 jours en avril.

Limites assumées : pas de saisonnalité, pas de prix d'hôtels intégré au score. Si on faisait ce projet en prod, on croiserait avec un dataset de réservations réelles.""")

add_footer(s, 4)


# ============================================================
# SLIDE 5 — AWS S3 + RDS
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
add_title_bar(s, "Stockage AWS — S3 Data Lake + RDS Data Warehouse",
              "Architecture cloud du Bloc 1")

# Bloc S3 à gauche
s3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4),
                        Inches(6.0), Inches(2.8))
s3.fill.solid(); s3.fill.fore_color.rgb = WHITE
s3.line.color.rgb = ORANGE; s3.line.width = Pt(2.0)
tf = s3.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "AWS S3  —  Data Lake"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ORANGE
for line in [
    "Bucket  s3://kayak-jedha-aymeric/raw/",
    "4 fichiers CSV bruts  ·  ~137 KB total",
    "    cities_coords.csv  (35 lignes)",
    "    weather_cities.csv  (35 lignes agrégées)",
    "    weather_daily_detail.csv  (175 lignes)",
    "    hotels_booking.csv  (679 lignes)",
    "Versioning ON  ·  Server-side encryption",
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13) if line.startswith("    ") else Pt(14)
    p.font.color.rgb = GRAY if line.startswith("    ") else NAVY
    p.space_before = Pt(2)

# Bloc RDS à droite
rds = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.83), Inches(1.4),
                         Inches(6.0), Inches(2.8))
rds.fill.solid(); rds.fill.fore_color.rgb = WHITE
rds.line.color.rgb = NAVY; rds.line.width = Pt(2.0)
tf = rds.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "AWS RDS PostgreSQL  —  Data Warehouse"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = NAVY
for line in [
    "Instance  db.t3.micro  ·  20 GB gp2",
    "Free tier 12 mois  ·  PostgreSQL 16",
    "2 tables nettoyées :",
    "    cities_weather  (35 lignes, PK city_id)",
    "    hotels  (679 lignes, FK city_id)",
    "Chiffrement SSL  ·  password fort",
    "Connexion testée depuis notebook 02",
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13) if line.startswith("    ") else Pt(14)
    p.font.color.rgb = GRAY if line.startswith("    ") else NAVY
    p.space_before = Pt(2)

# Bandeau coût et livrables en bas
add_text_box(s, Inches(0.5), Inches(4.5), Inches(12.5), Inches(2.5),
             ["Livrables vérifiés",
              "1 bucket S3 avec les CSV enrichis (météo + hôtels) par ville",
              "1 base SQL RDS où on retrouve ces mêmes données nettoyées",
              "2 cartes Plotly  :  map_destinations.html  +  map_hotels.html",
              "",
              "Coût mensuel  :  0 € (free tier 12 mois)  →  ≈ 13 €/mois après free tier"],
             size=14, bold_first=True, bullet=False, color=NAVY)

add_notes(s, """[~1 min — STOCKAGE AWS]

L'infrastructure AWS — c'est le cœur du Bloc 1, la partie « data engineering ».

Côté S3, mon Data Lake : un bucket kayak-jedha-aymeric, prefix raw/, avec 4 fichiers CSV pour 137 KB au total. C'est minuscule, mais l'architecture est la même qu'en prod, juste avec un volume bootcamp. J'ai activé le versioning et le chiffrement côté serveur.

Côté RDS, mon Data Warehouse : une instance db.t3.micro PostgreSQL 16, en free tier 12 mois. Deux tables : cities_weather avec les 35 villes et leur score, et hotels avec les 679 hôtels — clé étrangère sur city_id pour faire le lien.

Si le jury me demande pourquoi les Security Groups RDS sont ouverts à 0.0.0.0/0 : c'est un compromis bootcamp, pour pouvoir me connecter de n'importe quelle wifi. En prod, on restreindrait à des IP fixes ou un bastion.

Côté livrables : les 3 sont là — bucket S3, base RDS requêtable, et les 2 cartes Plotly.

Coût : 0 € pendant la formation grâce au free tier. Après, ça serait environ 13 €/mois — la RDS coûte plus que tout le reste.""")

add_footer(s, 5)


# ============================================================
# SLIDE 6 — Conclusion + perspectives
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
add_title_bar(s, "Conclusion & perspectives",
              "Ce que j'ai appris  ·  Ce que je ferais différemment en prod")

# Acquis
add_text_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.5),
             ["Acquis du Bloc 1"], size=20, bold_first=True, bullet=False, color=ORANGE)

acquis = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0),
                            Inches(6.0), Inches(4.5))
acquis.fill.solid(); acquis.fill.fore_color.rgb = WHITE
acquis.line.color.rgb = ORANGE; acquis.line.width = Pt(1.5)
tf = acquis.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
for i, line in enumerate([
    "Pipeline ETL bout-en-bout fonctionnel",
    "Scraping site dynamique (Selenium)",
    "Géocodage + appel API REST avec gestion d'erreurs",
    "Stockage cloud S3 + RDS PostgreSQL",
    "Modélisation relationnelle (PK / FK)",
    "Restitution Plotly interactive (cartes HTML)",
    "Reproductibilité  :  requirements.txt + .env.example + README",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = "✓  " + line
    p.font.size = Pt(14); p.font.color.rgb = NAVY
    p.space_before = Pt(6)

# Pour aller plus loin
add_text_box(s, Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.5),
             ["Pour aller en prod chez Kayak"], size=20, bold_first=True,
             bullet=False, color=NAVY)

prod = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0),
                          Inches(6.0), Inches(4.5))
prod.fill.solid(); prod.fill.fore_color.rgb = WHITE
prod.line.color.rgb = NAVY; prod.line.width = Pt(1.5)
tf = prod.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
for i, line in enumerate([
    "Orchestration  :  Airflow ou Step Functions (cron quotidien)",
    "Sécurité  :  IAM role EC2 → fini les credentials dans .env",
    "Robustesse  :  retry exponentiel, alerting CloudWatch",
    "Versioning  :  S3 versioning ON pour rollback rapide",
    "Tests  :  unitaires sur le parsing + end-to-end",
    "Score météo  :  poids appris depuis comportements réels",
    "Élargir le périmètre  :  Europe, puis monde",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = "→  " + line
    p.font.size = Pt(14); p.font.color.rgb = NAVY
    p.space_before = Pt(6)

# Merci
merci = s.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.5), Inches(0.5))
p = merci.text_frame.paragraphs[0]
p.text = "Merci pour votre attention  —  questions ?"
p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

add_notes(s, """[~30 s — CONCLUSION]

Pour conclure.

Ce que j'ai appris dans ce Bloc 1 : construire un pipeline ETL bout-en-bout, scraper un site dynamique, manipuler les APIs REST, stocker en cloud sur deux niveaux (lake et warehouse), modéliser une base relationnelle, et restituer en Plotly. Le tout reproductible avec requirements.txt et un README.

Ce que je ferais différemment pour aller en prod chez Kayak : d'abord orchestrer le pipeline avec Airflow pour qu'il tourne tous les jours automatiquement. Ensuite remplacer mes credentials .env par un IAM role attaché à une EC2. Ajouter du retry exponentiel, du monitoring CloudWatch. Apprendre les poids du score météo à partir de données réelles. Et bien sûr élargir le périmètre — 35 villes c'est un POC, en vrai Kayak couvre le monde.

Voilà, merci pour votre attention. Je suis prêt pour vos questions.""")

add_footer(s, 6)


# ============================================================
# SAVE
# ============================================================
out = "Presentation_Kayak.pptx"
prs.save(out)
print(f"OK -> {out}")
print(f"Slides : {len(prs.slides)}")
